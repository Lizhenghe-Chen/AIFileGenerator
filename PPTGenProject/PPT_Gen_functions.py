import json
import re
import random
from typing import Dict, List, Any, Optional
from pptx import Presentation
import os
from PPT_Prompt import get_ppt_generation_prompt
from openai import OpenAI
from config import OPENAI_CONFIG, PPT_CONFIG, PATHS, LOGGING_CONFIG


def generate_table_of_contents(slides_data: List[Dict[str, Any]]) -> Dict[str, Any]:
    """根据幻灯片内容生成目录"""
    toc_items = []

    for i, slide in enumerate(slides_data, 1):
        if slide.get("type") == "content":
            title = slide.get("title", f"第{i}部分")
            toc_items.append(title)

    return {
        "type": "content",
        "title": "目录",
        "content_type": "bullet_list",
        "content": toc_items,
        "has_image": False,
    }


def create_table_of_contents_slide(prs, toc_data: Dict[str, Any]):
    """创建目录幻灯片"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "目录"

    # 获取内容占位符并调整位置和大小
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear()

    # 添加目录项
    toc_items = toc_data.get("content", [])
    for i, item in enumerate(toc_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # 添加序号和内容
        p.text = f"{i + 1}. {item}"
        p.level = 0


def parse_content(content: str) -> Dict[str, Any]:
    """解析GPT返回的内容"""
    if LOGGING_CONFIG["show_debug_info"]:
        print("=" * 60)
        print("🤖 GPT 生成的原始内容:")
        print("=" * 60)
        print(content)
        print("=" * 60)

    try:
        # 尝试直接解析JSON
        parsed_data = json.loads(content)
        print("✅ JSON解析成功!")
        return parsed_data
    except json.JSONDecodeError:
        print("⚠️ 直接JSON解析失败，尝试提取JSON部分...")
        # 如果解析失败，尝试提取JSON部分
        json_match = re.search(r"\{.*\}", content, re.DOTALL)
        if json_match:
            try:
                parsed_data = json.loads(json_match.group())
                print("✅ 提取JSON解析成功!")
                return parsed_data
            except json.JSONDecodeError:
                print("❌ 提取JSON解析也失败")

        print("🔄 使用默认结构...")
        # 如果仍然失败，返回默认结构
        return {
            "title": "演示文稿",
            "slides": [{"type": "title", "title": "标题", "subtitle": "副标题"}],
        }


def create_title_slide(prs, slide_data: Dict[str, str]):
    """创建标题幻灯片"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = slide_data.get("title", "标题")
    subtitle.text = slide_data.get("subtitle", "副标题")


def get_available_content_layouts(prs) -> List[int]:
    """自动检测模板中可用的内容布局（必须同时包含标题和内容占位符的布局）"""
    available_layouts = []

    for i, layout in enumerate(prs.slide_layouts):
        # 跳过布局0（通常是标题页布局）
        if i == 0:
            continue

        # 检查是否有标题占位符（索引0）和内容占位符
        has_title = False
        has_content = False

        try:
            for placeholder in layout.placeholders:
                # 检查是否为标题占位符
                if placeholder.placeholder_format.idx == 0:
                    has_title = True
                # 检查是否为内容占位符（有text_frame且不是标题）
                elif (
                    hasattr(placeholder, "text_frame")
                    and placeholder.placeholder_format.idx != 0
                ):
                    has_content = True

            # 只有同时具备标题和内容占位符才认为是可用的内容布局
            if has_title and has_content:
                available_layouts.append(i)
                if LOGGING_CONFIG.get("show_debug_info", False):
                    print(f"  🔍 发现可用布局: {i} (标题✓ 内容✓)")
            else:
                if LOGGING_CONFIG.get("show_debug_info", False):
                    title_status = "✓" if has_title else "✗"
                    content_status = "✓" if has_content else "✗"
                    print(f"  ❌ 布局 {i} 不可用 (标题{title_status} 内容{content_status})")
        except Exception as e:
            if LOGGING_CONFIG.get("show_debug_info", False):
                print(f"  ⚠️ 检查布局 {i} 时出错: {e}")
            continue

    if LOGGING_CONFIG.get("show_debug_info", False):
        print(f"  📋 模板中可用的内容布局: {available_layouts}")

    return available_layouts


# 缓存布局检测结果，避免重复检测
_layout_cache = {}


def get_random_content_layout(prs) -> int:
    """随机选择一个内容布局"""
    # 检查是否使用自动检测
    use_auto_detection = PPT_CONFIG.get("auto_detect_layouts", True)

    if use_auto_detection:
        # 创建缓存键（基于模板的布局数量）
        cache_key = len(prs.slide_layouts)

        # 检查缓存
        if cache_key in _layout_cache:
            available_layouts = _layout_cache[cache_key]
            if LOGGING_CONFIG.get("show_debug_info", False):
                print(f"  📋 使用缓存的布局: {available_layouts}")
        else:
            # 自动检测可用布局并缓存
            available_layouts = get_available_content_layouts(prs)
            _layout_cache[cache_key] = available_layouts

        if not available_layouts:
            print("⚠️ 未检测到可用的内容布局，使用默认布局1")
            return 1
    else:
        # 使用配置文件中的布局列表
        available_layouts = PPT_CONFIG.get("available_content_layouts", [1])
        if not available_layouts:
            print("⚠️ 配置中无可用布局，使用默认布局1")
            return 1

    return random.choice(available_layouts)


def create_content_slide_with_layout(
    prs, slide_data: Dict[str, Any], layout_index: int
):
    """使用指定布局创建内容幻灯片"""
    content_type = slide_data.get("content_type", "bullet_list")

    # 使用指定的布局
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)

    # 设置标题
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get("title", "标题")

    # 根据布局类型处理内容
    content_data = slide_data.get("content", [])

    # 尝试找到内容占位符
    content_placeholder = None
    for shape in slide.placeholders:
        # 查找文本占位符（通常是索引1或其他）
        if (
            hasattr(shape, "text_frame") and shape.placeholder_format.idx != 0
        ):  # 0通常是标题
            content_placeholder = shape
            break

    if content_placeholder is None:
        # 如果没有找到占位符，尝试使用索引1
        try:
            content_placeholder = slide.placeholders[1]
        except (IndexError, KeyError):
            print(f"⚠️ 布局 {layout_index} 没有可用的内容占位符，跳过内容添加")
            return

    # 处理文本内容
    if hasattr(content_placeholder, "text_frame"):
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        if content_type == "bullet_list" and isinstance(content_data, list):
            # 项目符号列表
            for i, item in enumerate(content_data):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = item
                p.level = 0

        elif content_type == "paragraph":
            # 段落文字
            p = text_frame.paragraphs[0]
            p.text = str(content_data)

        elif content_type == "title_paragraph" and isinstance(content_data, dict):
            # 小标题加段落
            subtitle = content_data.get("subtitle", "")
            text = content_data.get("text", "")

            # 小标题
            p1 = text_frame.paragraphs[0]
            p1.text = subtitle

            # 段落文字
            p2 = text_frame.add_paragraph()
            p2.text = text

        else:
            # 默认处理：转换为字符串
            p = text_frame.paragraphs[0]
            p.text = str(content_data)


def create_content_slide(prs, slide_data: Dict[str, Any]):
    """创建内容幻灯片（支持随机布局）"""
    # 检查是否启用随机布局
    use_random_layouts = PPT_CONFIG.get("use_random_layouts", True)

    if use_random_layouts:
        # 随机选择布局
        layout_index = get_random_content_layout(prs)

        if LOGGING_CONFIG.get("show_debug_info", False):
            print(f"  🎲 随机选择布局: {layout_index}")
    else:
        # 使用默认布局1
        layout_index = 1
        if LOGGING_CONFIG.get("show_debug_info", False):
            print(f"  📄 使用默认布局: {layout_index}")

    try:
        create_content_slide_with_layout(prs, slide_data, layout_index)
    except Exception as e:
        print(f"⚠️ 使用布局 {layout_index} 失败: {e}")
        print("🔄 回退到默认布局 1")
        # 如果随机布局失败，回退到布局1
        create_content_slide_with_layout(prs, slide_data, 1)


def get_template_path(design_number: Optional[int]) -> str:
    """获取模板文件路径"""
    return PATHS["template_path_format"].format(design_number)


def create_presentation(
    ppt_data: Dict[str, Any],
    design_number: Optional[int] = None,
    custom_filename: Optional[str] = None,
) -> str:
    """创建完整的演示文稿"""
    if design_number is None:
        design_number = PPT_CONFIG["default_design_number"]

    template_path = get_template_path(design_number)

    # 使用模板文件创建演示文稿
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        print(f"🎨 使用模板: {template_path}")
    else:
        prs = Presentation()
        print(f"⚠️ 模板文件不存在，使用默认模板: {template_path}")

    # 获取演示文稿标题和建议文件名
    presentation_title = ppt_data.get("title", "演示文稿")
    suggested_filename = ppt_data.get("filename", "presentation")

    # 确定最终文件名
    if custom_filename:
        filename = f"{custom_filename}.pptx"
    else:
        filename = f"{suggested_filename}.pptx"

    # 创建Output文件夹（如果不存在）
    output_dir = "Output"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        print(f"📁 创建输出目录: {output_dir}")

    # 生成完整的文件路径
    full_path = os.path.join(output_dir, filename)

    print(f"\n📊 开始创建演示文稿: {presentation_title}")
    print(f"📁 文件路径: {full_path}")

    slides = ppt_data.get("slides", [])

    # 生成目录数据
    toc_data = generate_table_of_contents(slides)
    print(f"📋 已生成目录，包含 {len(toc_data['content'])} 个章节")

    # 创建幻灯片
    slide_counter = 0

    for i, slide_data in enumerate(slides):
        slide_type = slide_data.get("type", "content")
        slide_title = slide_data.get("title", f"幻灯片 {i + 1}")

        slide_counter += 1
        print(f"  📄 创建第 {slide_counter} 页: {slide_title} ({slide_type})")

        if slide_type == "title":
            create_title_slide(prs, slide_data)
            # 在标题页后插入目录页
            if len(toc_data["content"]) > 0:
                slide_counter += 1
                print(f"  📋 创建第 {slide_counter} 页: 目录 (table_of_contents)")
                create_table_of_contents_slide(prs, toc_data)
        else:
            create_content_slide(prs, slide_data)

    # 保存文件
    prs.save(full_path)
    return full_path


def get_openai_client(base_url: Optional[str] = None, api_key: Optional[str] = None) -> OpenAI:
    """创建OpenAI客户端"""
    if base_url is None:
        base_url = OPENAI_CONFIG["base_url"]
    if api_key is None:
        api_key = OPENAI_CONFIG["api_key"]
    return OpenAI(base_url=base_url, api_key=api_key)


def generate_ppt_content(
    user_input: str,
    expected_slides: Optional[int] = None,
    base_url: Optional[str] = None,
    api_key: Optional[str] = None,
    model_path: Optional[str] = None,
) -> str:
    """使用GPT根据用户输入生成PPT内容"""
    if expected_slides is None:
        expected_slides = PPT_CONFIG["default_expected_slides"]
    if base_url is None:
        base_url = OPENAI_CONFIG["base_url"]
    if api_key is None:
        api_key = OPENAI_CONFIG["api_key"]
    if model_path is None:
        model_path = OPENAI_CONFIG["model_path"]

    client = get_openai_client(base_url, api_key)
    prompt = get_ppt_generation_prompt(user_input, expected_slides)

    messages = [{"role": "user", "content": prompt}]

    response = client.chat.completions.create(
        model=model_path,
        messages=messages,
        stream=False,
    )

    return response.choices[0].message.content


def generate_ppt_from_user_input(
    user_input: str,
    expected_slides: Optional[int] = None,
    custom_filename: Optional[str] = None,
    design_number: Optional[int] = None,
    base_url: Optional[str] = None,
    api_key: Optional[str] = None,
    model_path: Optional[str] = None,
) -> str:
    """根据用户输入生成PPT的完整流程"""
    # 使用配置文件的默认值
    if expected_slides is None:
        expected_slides = PPT_CONFIG["default_expected_slides"]
    if design_number is None:
        design_number = PPT_CONFIG["default_design_number"]
    if base_url is None:
        base_url = OPENAI_CONFIG["base_url"]
    if api_key is None:
        api_key = OPENAI_CONFIG["api_key"]
    if model_path is None:
        model_path = OPENAI_CONFIG["model_path"]

    if LOGGING_CONFIG["show_progress"]:
        print(f"🚀 正在根据用户需求生成PPT...")
        print(f"📝 用户需求: {user_input}")
        print(f"📊 期望页数: {expected_slides}页")
        print(f"🎨 使用模板: Design-{design_number}.pptx")

    # 生成内容
    content = generate_ppt_content(
        user_input, expected_slides, base_url, api_key, model_path
    )
    print("✅ GPT内容生成完成！")

    # 解析内容
    ppt_data = parse_content(content)
    print("✅ 内容解析完成！")

    # 创建PPT
    saved_filename = create_presentation(ppt_data, design_number, custom_filename)
    print(f"✅ PPT文件已创建：{saved_filename}")

    return saved_filename
