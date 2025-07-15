from pptx import Presentation
from pptx.util import Inches, Pt

# 创建一个新的演示文稿
prs = Presentation()

# 添加一个标题幻灯片
slide_layout = prs.slide_layouts[0]  # 0 是标题幻灯片布局
slide = prs.slides.add_slide(slide_layout)

# 设置标题和副标题
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "你好，Python PPTX"
subtitle.text = "这是用 Python 创建的第一个幻灯片"

# 添加一个内容幻灯片
slide_layout = prs.slide_layouts[1]  # 1 是标题和内容布局
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
content = slide.placeholders[1]

title.text = "第二页标题"
content.text = "这是内容\n- 项目1\n- 项目2\n- 项目3"

# 保存文件
prs.save("简单演示文稿.pptx")
print("PPTX 文件已创建：简单演示文稿.pptx")