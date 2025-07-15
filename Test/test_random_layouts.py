#!/usr/bin/env python3
"""
测试随机布局功能
"""
import os
from PPT_Gen_functions import (
    get_random_content_layout,
    get_available_content_layouts,
    create_presentation,
    generate_ppt_from_user_input,
)
from pptx import Presentation
from config import PPT_CONFIG


def test_auto_layout_detection():
    """测试自动布局检测功能"""
    print("🧪 测试自动布局检测功能...")
    
    # 使用默认模板测试
    prs = Presentation()
    available_layouts = get_available_content_layouts(prs)
    print(f"默认模板检测到的布局: {available_layouts}")
    
    # 测试使用Design模板
    design_templates = [1, 2, 3, 4, 5, 6, 7]
    for design_num in design_templates:
        template_path = f"Designs/Design-{design_num}.pptx"
        if os.path.exists(template_path):
            try:
                prs = Presentation(template_path)
                available_layouts = get_available_content_layouts(prs)
                print(f"Design-{design_num} 检测到的布局: {available_layouts}")
            except Exception as e:
                print(f"⚠️ 无法加载 Design-{design_num}: {e}")
        else:
            print(f"⚠️ Design-{design_num} 模板文件不存在")


def test_random_layout_selection():
    """测试随机布局选择功能"""
    print("🧪 测试随机布局选择功能...")

    # 创建一个空的演示文稿来测试
    prs = Presentation()

    # 测试多次随机选择
    layouts = []
    for i in range(10):
        layout = get_random_content_layout(prs)
        layouts.append(layout)
        print(f"第 {i+1} 次选择: 布局 {layout}")

    print(f"\n📊 布局选择统计:")
    print(f"选择的布局: {sorted(set(layouts))}")
    print(f"可用布局: {PPT_CONFIG['available_content_layouts']}")
    print(
        f"所有选择都在可用范围内: {all(layout in PPT_CONFIG['available_content_layouts'] for layout in layouts)}"
    )


def test_ppt_generation_with_random_layouts():
    """测试使用随机布局生成PPT"""
    print("\n🧪 测试随机布局PPT生成...")

    # 测试数据
    test_input = (
        "创建一个关于人工智能发展的演示文稿，包括AI历史、当前应用、未来趋势等内容"
    )

    try:
        # 生成PPT
        filename = generate_ppt_from_user_input(
            user_input=test_input,
            expected_slides=5,
            custom_filename="test_random_layouts",
            design_number=3,
        )
        print(f"✅ 测试PPT生成成功: {filename}")
        return True
    except Exception as e:
        print(f"❌ 测试PPT生成失败: {e}")
        return False


def main():
    """主测试函数"""
    print("🚀 开始测试随机布局功能")
    print("=" * 50)

    # 显示当前配置
    print("📋 当前配置:")
    print(f"  使用随机布局: {PPT_CONFIG.get('use_random_layouts', True)}")
    print(f"  自动检测布局: {PPT_CONFIG.get('auto_detect_layouts', True)}")
    print(f"  手动配置布局: {PPT_CONFIG.get('available_content_layouts', [])}")
    print("=" * 50)

    # 测试自动布局检测
    test_auto_layout_detection()
    
    print("\n" + "=" * 50)
    
    # 测试随机布局选择
    test_random_layout_selection()

    # 测试PPT生成（需要有效的OpenAI配置）
    print("\n" + "=" * 50)
    print("⚠️ 注意：以下测试需要有效的OpenAI API配置")
    user_choice = input("是否继续测试PPT生成？(y/n): ").lower().strip()

    if user_choice == "y":
        success = test_ppt_generation_with_random_layouts()
        if success:
            print("\n✅ 所有测试通过！")
        else:
            print("\n❌ PPT生成测试失败")
    else:
        print("\n⏭️ 跳过PPT生成测试")

    print("\n🎉 测试完成！")


if __name__ == "__main__":
    main()
